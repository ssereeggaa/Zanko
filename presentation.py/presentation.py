from pptx import Presentation
from pptx.util import Inches

# Створення презентації
prs = Presentation()

def add_slide(title, content, image_path=None):
    slide_layout = prs.slide_layouts[5]  # Заголовок і текст
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[0]
    title_placeholder.text = title
    content_placeholder.text = content
    
    if image_path:
        slide.shapes.add_picture(image_path, Inches(1), Inches(3), width=Inches(5))

# Додавання титульного слайду
slide_layout = prs.slide_layouts[0]  # Титульний слайд
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Системи машинної геометрії"
slide.placeholders[1].text = "Група: [Ваша група]\nПрізвище та ім'я: [Ваші дані]"

# Додавання змісту
slides_content = [
    ("Вступ", "Машинна геометрія – це розділ комп'ютерної графіки, який допомагає працювати з формами."),
    ("Основні поняття", "У машинній геометрії використовуються прості фігури: точки, лінії, криві, поверхні."),
    ("Види систем", "Системи поділяються на проєктування, виробництво та аналіз."),
    ("Системи проєктування", "Ці системи допомагають створювати креслення і 3D моделі."),
    ("Системи виробництва", "Вони використовуються для керування верстатами."),
    ("Системи аналізу", "Допомагають перевіряти міцність і стійкість моделей."),
    ("Математичні основи", "Використовуються рівняння для роботи з формами та об'єктами."),
    ("Геометричні перетворення", "Можна змінювати розмір, повертати і переміщувати об'єкти."),
    ("Де це використовується?", "Це важливо у будівництві, дизайні, іграх, медицині."),
    ("Висновки", "Системи машинної геометрії важливі для технологічного розвитку.")
]

for title, content in slides_content:
    add_slide(title, content)

# Збереження презентації
prs.save("machine_geometry_presentation.pptx")
print("Презентація збережена як machine_geometry_presentation.pptx")
